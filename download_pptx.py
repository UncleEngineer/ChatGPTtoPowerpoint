# !pip install python-pptx

from pptx import Presentation
from pptx.util import Inches
from google.colab import files

# ข้อมูล JSON ที่ได้รับ
jsondata = {
  "title": "การปลูกผักในบ้าน",
  "content": [
    {
      "heading": "1. เลือกผักที่เหมาะสม",
      "details": [
        "ผักที่ปลูกง่ายในบ้าน: ผักสลัด, ผักบุ้ง, ต้นหอม, โหระพา",
        "เลือกผักที่สามารถเติบโตได้ในที่ที่มีแสงน้อย"
      ]
    },
    {
      "heading": "2. การให้แสง",
      "details": [
        "ผักต้องการแสงแดดพอประมาณ",
        "ใช้ไฟ LED (grow light) สำหรับบ้านที่มีแสงไม่พอ"
      ]
    },
    {
      "heading": "3. กระถางและดิน",
      "details": [
        "เลือกกระถางที่มีรูระบายน้ำ",
        "ใช้ดินที่มีคุณสมบัติระบายน้ำได้ดี"
      ]
    },
    {
      "heading": "4. การรดน้ำ",
      "details": [
        "รดน้ำอย่างพอเหมาะ ไม่มากเกินไป"
      ]
    },
    {
      "heading": "5. การดูแลและเก็บเกี่ยว",
      "details": [
        "สังเกตการเจริญเติบโตของผัก",
        "เก็บเกี่ยวเมื่อผักโตเต็มที่"
      ]
    },
    {
      "heading": "6. การเลือกสถานที่สำหรับปลูกผักในบ้าน",
      "details": [
        "เลือกสถานที่ที่มีแสงสว่าง เช่น หน้าต่างที่หันไปทางทิศใต้",
        "พิจารณาความสะดวกในการดูแล เช่น สถานที่ที่สามารถเข้าถึงได้ง่าย"
      ]
    },
    {
      "heading": "7. การใช้กระถางและวัสดุปลูก",
      "details": [
        "เลือกกระถางที่มีรูระบายน้ำ",
        "ใช้วัสดุปลูกที่ระบายน้ำได้ดี เช่น ดินผสมทราย"
      ]
    },
    {
      "heading": "8. การรดน้ำผัก",
      "details": [
        "รดน้ำอย่างสม่ำเสมอเมื่อดินเริ่มแห้ง",
        "หลีกเลี่ยงการรดน้ำบนใบเพื่อป้องกันโรค"
      ]
    },
    {
      "heading": "9. การใส่ปุ๋ย",
      "details": [
        "ใช้ปุ๋ยอินทรีย์หรือปุ๋ยคอกที่มีธาตุอาหารครบถ้วน",
        "ใช้ปุ๋ยน้ำสำหรับผักเพื่อเพิ่มการเจริญเติบโต"
      ]
    },
    {
      "heading": "10. การควบคุมอุณหภูมิในบ้าน",
      "details": [
        "รักษาอุณหภูมิระหว่าง 18-24 องศาเซลเซียส",
        "หลีกเลี่ยงการวางใกล้แหล่งความร้อน"
      ]
    },
    {
      "heading": "11. การดูแลหลังการปลูก",
      "details": [
        "ตัดแต่งใบที่แห้งหรือเสียหายออก",
        "สามารถขยายพันธุ์จากการปักชำหรือเมล็ด"
      ]
    },
    {
      "heading": "12. การควบคุมโรคและแมลง",
      "details": [
        "สังเกตอาการของพืช เช่น ใบเหลืองหรือรากเน่า",
        "ใช้ยาฆ่าแมลงอินทรีย์ เช่น น้ำส้มสายชูผสมน้ำ"
      ]
    },
    {
      "heading": "13. การเก็บเกี่ยวผัก",
      "details": [
        "เก็บผักเมื่อถึงเวลาที่เหมาะสม เช่น ผักสลัด",
        "เก็บผักในช่วงเช้าเพื่อรักษาความสดใหม่"
      ]
    },
    {
      "heading": "14. การปลูกผักในระบบน้ำ (Hydroponics)",
      "details": [
        "ใช้ระบบ hydroponics ที่ไม่ใช้ดินแต่ใช้สารละลายธาตุอาหาร",
        "ประหยัดน้ำและพื้นที่"
      ]
    },
    {
      "heading": "15. ประโยชน์ของการปลูกผักในบ้าน",
      "details": [
        "การประหยัดค่าใช้จ่ายและลดค่าใช้จ่ายในครัวเรือน",
        "อาหารสดและปลอดภัยจากสารเคมี",
        "เป็นกิจกรรมที่ผ่อนคลายและลดความเครียด"
      ]
    }


  ]
}

# สร้างไฟล์ PowerPoint
prs = Presentation()

# สไลด์แรก: ชื่อเรื่อง
slide_0 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_0.shapes.title
title.text = jsondata["title"]

# สร้างสไลด์สำหรับเนื้อหา
for section in jsondata["content"]:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = section["heading"]
    
    content = "\n".join(section["details"])
    textbox = slide.shapes.placeholders[1]
    textbox.text = content

# บันทึกไฟล์ PowerPoint
output_filename = "/content/การปลูกผักในบ้าน2.pptx"
prs.save(output_filename)

# ให้ผู้ใช้ดาวน์โหลดไฟล์
files.download(output_filename)
